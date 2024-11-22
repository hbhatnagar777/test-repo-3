# -*- coding: utf-8 -*-

# --------------------------------------------------------------------------
# Copyright Commvault Systems, Inc.
# See LICENSE.txt in the project root for
# license information.
# --------------------------------------------------------------------------

"""Main module for calling wikipedia API to create mail body

Wiki class is defined in this module.

Wiki: Calls wikipedia API and fetches random pages.

Wiki:
    __init__(log)   --  Initializes the wiki object

    _set_params()   --  Sets parameter dict for calling respective wikipedia API

    create_message_from_wiki()  --  Fetches page title and content of a random page from wikipedia

    __get_wikiapi_service() --  Returns a random language wiki api service.
    Useful to create unicode mails.

    make_request()  --  This method make a request to API endpoint.

    create_docx()   --  This method creates a new documents and saves it into the directory
    defined in constants file.

    make_valid_filename()  --  This method returns the valid file name for Windows OS

"""

from __future__ import unicode_literals
import os
import time
import shutil
from random import randrange
import unicodedata as ud
import codecs
import requests
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import cm
from reportlab.pdfgen import canvas
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from .exception import CVCloudException
from . import constants
import csv
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


class Wiki:
    """Class to interact with wikipedia API to fetch random articles"""

    def __init__(self, log):
        """Initializes the wiki object

                Args:

                    log (object) -- instance of the logger module

                Returns:

                    object  --  Instance of Wiki class

        """
        self.log = log
        self.params = {'action': 'query', 'format': 'json'}
        self.app_name = self.__class__.__name__

        _all_unicode = ''.join(chr(i) for i in range(65536))
        self._unicode_letters = ''.join(c for c in _all_unicode
                                        if ud.category(c) == 'Lu' or ud.category(c) == 'Ll')
        self.unicode_letters = set(self._unicode_letters)

    def _set_params(self, content_type, wiki_id=None):
        """Sets parameter dict for calling respective wikipedia API

                Args:

                    content_type (string) -- Type of the content to fetch from wikipedia

                                            Allowed values:

                                                subject -- This returns the title of the page

                                                body -- This returns the content of the page

                    wiki_id (int)  -- Page id of the page got from wikipedia API

                Returns:

                    dict -- Dict with required parameters to call wiki API

        """
        self.log.info('setting parameters for wiki API..')

        if content_type == 'subject':

            add_params = {
                'list': 'random',
                'rnnamespace': 0}

        elif content_type == 'body':

            add_params = {
                #'prop': 'revisions',
                'prop': 'extracts',
                #'rvprop': 'content',
                'explaintext': 1,
                'exsectionformat': 'plain',
                'pageids': int(wiki_id)}
        else:
            raise CVCloudException(self.app_name, '104')

        add_params.update(self.params)
        return add_params

    def create_message_from_wiki(self, unicode_data=False):
        """Fetches page title and content of a random page from wikipedia

                Args:

                    unicode_data (boolean)  --  True if unicode data has to be fetched

                                                Default: False

                Returns:

                    Dict of following type: {'subject': <some subject>, 'body': <page content>}

        """
        message = {}
        try:
            self.log.info('Fetching random article id from wiki')
            wiki_api_service = self.__get_wikiapi_service(unicode_data)
            params_dict = self._set_params('subject')
            self.log.info('params dict for wiki API:%s', params_dict)
            response = self.make_request(
                req_url=wiki_api_service, req_body=params_dict)
            message['subject'] = response.get(
                'query').get('random')[0].get('title')
            self.log.info(u'message subject: %s', message['subject'])
            wiki_id = response.get('query').get('random')[0].get('id')
            params_dict = self._set_params('body', wiki_id)
            self.log.info('params dict for wiki API:%s', params_dict)
            response = self.make_request(
                req_url=wiki_api_service, req_body=params_dict)
            # message['body'] = response.get('query').get('pages').get(
            #    str(wiki_id)).get('revisions')[0].get('*')
            message['body'] = response.get('query').get('pages').get(
                str(wiki_id)).get('extract')
            self.log.info('Message content got from wiki')
            return message
        except Exception as excp:
            self.log.exception('Exception while creating message from wiki')
            raise CVCloudException(self.app_name, '103', str(excp))

    def __get_wikiapi_service(self, unicode_data=False):
        """Returns a random language wiki api service. Useful to create unicode mails.

                Args:

                    unicode_data (boolean):

                        Default: False, if true it returns a wiki service other than English

                Returns:

                    wiki_api_service (str)  --  Wikipedia API endpoint

        """

        base_url = 'https://{0}.wikipedia.org/w/api.php?'
        wiki_language_list = [
            'en',
            'ru',
            'vi',
            'ja',
            'zh',
            'uk',
            'fa',
            'ar',
            'sh',
            'ko',
            'bg',
            'kk',
            'sk',
            'hy',
            'he',
            'ce',
            'el',
            'be',
            'hi',
            'th',
            'ur',
            'ka',
            'ta',
            'mk',
            'ne',
            'tg',
            'te',
            'ky',
            'zh-yue',
            'bn',
            'ml',
            'mr',
            'cv',
            'ba',
            'gu',
            'pa',
            'si',
            'or',
            'sa',
            'mrj',
            'os',
            'ilo']
        if not unicode_data:
            self.log.info(
                'Unicode data is set to False. Will fetch content from english wiki api')
            url = base_url.format('en')
        else:
            self.log.info(
                'Unicode data is set to True. Will fetch content from unicode wiki api')
            random_number = randrange(len(wiki_language_list))
            url = base_url.format(wiki_language_list[random_number])
        self.log.info('will fetch wiki api: %s', url)
        return url

    def make_request(self, req_type='GET', req_url=None, req_body=None):
        """This method make a request to API endpoint.

                Args:

                    req_type (str) -- Type of HTTP request

                                        Allowed values:
                                            GET
                                            POST

                    req_url (str)  --  URL to make request to

                    req_body (dict) -- URL parameters or POST data

                Returns:

                    responce (json)

        """
        response = None
        try:
            headers = {
                'Accept': 'application/json',
                'Content-Type': 'application/json'}
            if req_type.upper() == 'GET':
                response = requests.get(
                    req_url, params=req_body, headers=headers)
            if req_type.upper() == 'POST':
                response = requests.post(
                    req_url, data=req_body, headers=headers)
            self.log.info('URL requested: %s', response.url)
            if response.ok and response.status_code == 200:
                self.log.info(
                    'Got HTTP response OK. Try to convert to JSON decoded response.')
                json_decoded = response.json()
                return json_decoded
            else:
                self.log.exception('Error in wiki response.')
                raise CVCloudException(
                    self.app_name, '101', str(
                        response.raise_for_status()))
        except Exception as ex:
            self.log.exception('Exception while making HTTP request')
            raise CVCloudException(self.app_name, '102', str(ex))

    def create_docx(
            self,
            no_of_docs=5,
            unicode_data=False,
            word=True,
            pdf=False,
            google_docs=False,
            csvfile=False,
            pptx=False,
            xlsx=False,
            image=False,
            html=False,
            code=False,
            folder=False,
            folder_name=None,
            ca_type='gdrive'):
        """This method creates a new docx document and saves it into the directory
            defined in constants file.

                Args:

                    no_of_docs (int)  --  Number of docx documents to create

                    unicode_data (boolean)  --  True if unicode document has to be created

                                                Default: False

                    word  (boolean)  --  True if word files have to be created

                    pdf  (boolean)  --  True if pdf files need to be created

                    google_docs  (boolean)  --  True if google docs files need to be created

                    csv  (boolean)  --  True if csv files need to be created

                    pptx  (boolean)  --  True if power point files need to be created

                    xlsx  (boolean)  --  True if excel files need to be created

                    image  (boolean)  --  True if image files need to be created

                    html  (boolean)  --  True if html files need to be created

                    code  (boolean)  --  True if source code files need to be created

                    folder  (boolean)  --  True if folder needs to be created

                    folder_name    (str)   --  Path where the files need to be saved

                    ca_type  (str)  --  Cloud Apps instance type
                                            Valid Values:
                                            gmail

                                            gdrive

                                            onedrive

        """
        try:
            if ca_type == 'gdrive':
                doc_path = constants.GDRIVE_DOCUMENT_DIRECTORY
                if folder_name:
                    doc_path = os.path.join(doc_path, folder_name)
            elif ca_type == 'gmail':
                doc_path = constants.GMAIL_DOCUMENT_DIRECTORY
                if folder_name:
                    doc_path = os.path.join(doc_path, folder_name)
            elif ca_type == 'onedrive':
                doc_path = constants.ONEDRIVE_DOCUMENT_DIRECTORY
                if folder_name:
                    doc_path = os.path.join(doc_path, folder_name)
            else:
                self.log.error('Invalid ca_type')
                raise CVCloudException(self.app_name, '106')

            if os.path.exists(doc_path):
                time.sleep(5)
                shutil.rmtree(doc_path)

            os.makedirs(doc_path)

            for i in range(no_of_docs):
                time.sleep(1) # to ensure only 1 file is created per second
                if word:
                    document = Document()
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    document.add_heading(file_name)
                    paragraph = document.add_paragraph(wiki_message['body'])
                    doc_format = paragraph.paragraph_format
                    doc_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                    filename = os.path.join(
                        doc_path,
                        f'{file_name}.docx')
                    document.save(filename)

                if pdf:
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.pdf')

                    can = canvas.Canvas(out_file, pagesize=letter)
                    width, height = letter
                    text_object = can.beginText()
                    text_object.setTextOrigin(width - 19 * cm, height - 2.5 * cm)
                    text_object.setFont("Helvetica-Oblique", 14)
                    text_object.textLines(wiki_message['body'])
                    can.drawText(text_object)
                    can.showPage()
                    can.save()

                if google_docs:
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.txt')
                    with codecs.open(out_file, 'w', 'utf-8') as output:
                        output.write(wiki_message['body'])
                    output.close()

                if csvfile:
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    data = [wiki_message['body'].split("\n")]
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.csv')
                    with open(out_file, "w", newline="", encoding="utf-8") as csvFile:
                        writer = csv.writer(csvFile)
                        writer.writerows(data)
                    csvFile.close()
                if pptx:
                    prs = Presentation()
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.pptx')
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    text_box = slide.shapes.add_textbox(left=0, top=0, width=Inches(10), height=Inches(7))
                    text_frame = text_box.text_frame
                    text_frame.text = wiki_message['body']
                    prs.save(out_file)

                if xlsx:
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.xlsx')
                    df = pd.DataFrame({'Data': [wiki_message['body']]})
                    df.to_excel(out_file, index=False, header=False)

                if image:
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.jpg')
                    img = Image.new('RGB', (500, 500), color='white')
                    draw = ImageDraw.Draw(img)
                    font = ImageFont.load_default()
                    draw.text((10, 10), wiki_message['body'], fill='black', font=font)
                    img.save(out_file)

                if html:
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.html')
                    html_content = f"<html><head><title>Example HTML</title></head><body>{wiki_message['body']}</body></html>"
                    with codecs.open(out_file, "w", encoding="utf-8") as html_file:
                        html_file.write(html_content)
                    html_file.close()

                if code:
                    wiki_message = self.create_message_from_wiki(unicode_data)
                    file_name = self.make_valid_filename(wiki_message['subject'])
                    out_file = os.path.join(
                        doc_path,
                        f'{file_name}.c')

                    c_code = f"""
                    #include <stdio.h>

                    int main() {{
                        printf("{wiki_message['body']}\\n");
                        return 0;
                    }}
                    """
                    with open(out_file, "w", encoding="utf-8") as c_file:
                        c_file.write(c_code)

            if folder:
                wiki_message = self.create_message_from_wiki(unicode_data)
                folder_name = self.make_valid_filename(wiki_message['subject'])
                folder_path = os.path.join(doc_path, folder_name)
                os.makedirs(folder_path)
                self.create_docx(
                    no_of_docs=1,
                    unicode_data=unicode_data,
                    word=word,
                    pdf=pdf,
                    google_docs=google_docs,
                    csvfile=csvfile,
                    pptx=pptx,
                    xlsx=xlsx,
                    image=image,
                    html=html,
                    code=code,
                    folder_name=folder_name,
                    ca_type=ca_type)

        except Exception as excp:
            self.log.exception('Exception while creating documents.')
            raise CVCloudException(self.app_name, '105', str(excp))

    def make_valid_filename(self, file_name):
        """This method returns the valid file name for Windows OS
        after removing invalid characters.

                Args:

                    file_name  (str)  --  Name of the file

                Returns:

                    file_name  (str)  --  Returns a valid filename

        """
        #valid_chars = "-_.() %s%s%s" % (string.ascii_letters, string.digits, self.unicode_letters)
        #file_name = ''.join(c for c in file_name if c in valid_chars)
        file_name = f'file{str(int(time.time()))}'
        self.log.info('Valid filename for the file: %s', file_name)
        return file_name
