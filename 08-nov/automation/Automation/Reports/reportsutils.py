# -*- coding: utf-8 -*-

# ---------------------------------------------------------------------------
# Copyright Commvault Systems, Inc.
# See LICENSE.txt in the project root for
# license information.
# ---------------------------------------------------------------------------
""" This module has utility functions.

    get_reports_config()    --      Returns reports config file path.

    get_commcell_name()     --      Gets commcell name from config file.

    is_custom_report()      --      Checks if URL is custom report's URL.

"""

import os
import docx
from pptx import Presentation
from AutomationUtils.config import get_config
from AutomationUtils.constants import AUTOMATION_DIRECTORY
from AutomationUtils import config
from AutomationUtils.machine import Machine
from py_pdf_parser import tables
from py_pdf_parser.loaders import load_file

CONSTANTS = config.get_config()


def get_reports_config():
    """
    Returns reports config file path
    """
    return get_config(
        json_path=os.path.join(
            AUTOMATION_DIRECTORY, "Reports", "REPORTS_CONFIG.json"
        )
    )


from datetime import date, timedelta


def get_startdt_string(interval_type):
    """
        Get expected date based on 'type',

        Args:
            interval_type (int): type should be selected from DAILY_TYPE, WEEKLY_TYPE
        MONTHLY_TYPE from declaration

        Returns (String): returns expected string  depending on type
        """
    #  if today is '2020-03-06 00:00:00.0', this should return '2020-03-05 00:00:00.0'
    if interval_type == int(1):  # return previous day date
        return str(date.today() - timedelta(1)) + " 00:00:00.0"

    #  if today is '2020-03-06 00:00:00.0', this should return '2020-02-01 00:00:00.0'
    if interval_type == int(2):  # return previous month dat
        last_day_of_prev_month = date.today().replace(day=1) - timedelta(days=1)
        start_day_of_prev_month = date.today().replace(day=1) - timedelta(
            days=last_day_of_prev_month.day)
        return str(start_day_of_prev_month) + " 00:00:00.0"

    #  if today is '2020-03-06 00:00:00.0', this should return '2020-02-24 00:00:00.0'
    if interval_type == int(3):  # return previous week date
        today = date.today()
        return str(today - timedelta(days=today.weekday(), weeks=1)) + " 00:00:00.0"
    if interval_type == int(4):  # return range of last completed week
        today = date.today()
        return [str(today - timedelta(days=today.weekday(), weeks=1)) + " 00:00:00.0",
                str(today - timedelta(days=today.weekday(), weeks=2)) + " 00:00:00.0"]


def get_commcell_name(commcell_object):
    """
    Gets commcell name from config file, if its blank returns commcell name from commcell object
    Args:
        commcell_object: commcell object
    Returns:commcell name
    """
    if CONSTANTS.Reports.CommCellName_reports == "":
        return commcell_object.commserv_name
    return CONSTANTS.Reports.CommCellName_reports


def is_custom_report(url):
    """check if url is custom report's url"""
    if 'reportsplus' in url:
        return True
    return False


class PPTManager:
    """Provides operations on pptx files"""

    def __init__(self, file):
        self.ppt_file = file
        self.presentation = Presentation(self.ppt_file)

    def get_number_of_slides(self):
        """
        Read number of slide from present in ppt file
        Returns                       (int)      --        number of slides
        """
        return len(self.presentation.slides)

    def get_text_from_slide(self, slide_number):
        """
        Get text from specified slide
        Args:
            slide_number           (int)        --         Slide number
        Returns                    (String)     --         text present in slide
        """
        _text = []
        _slide = self.presentation.slides[slide_number]
        for shape in _slide.shapes:
            if shape.has_text_frame and shape.text != "":
                _text.append(shape.text)
        return _text

    def get_table_data(self, slide_number):
        """
        Get table data from specified slide number
        Args:
            slide_number            (int)      --          Slide number
        Returns                     (Dict)     --          Dictionary of table content
        """
        _slide = self.presentation.slides[slide_number - 1]
        _table = None
        for shape in _slide.shapes:
            if shape.has_table:
                _table = shape.table  # As of now doing for 1 table only
                break
        if not _table:  # If table is not there then no need to continue further
            raise Exception("Table is not present in [%s] slide" % slide_number)
        headings = []
        #  Collect headings of the table
        for each_col_count in range(0, len(list(_table.columns))):
            headings.append(_table.cell(0, each_col_count).text)
        #  Create dictionary of column name as key and append column values
        table_text = {}
        heading_index = 0
        for c_index in range(0, len(list(_table.columns))):
            temp = []
            for r_index in range(1, len(list(_table.rows))):
                temp.append(_table.cell(r_index, c_index).text)
            table_text.update({headings[heading_index]: temp})
            heading_index += 1
        return table_text


class CSUtils:
    """Reports folder related utils in CommServe"""

    def __init__(self, commcell_obj):
        """
        Args:
            commcell_obj: object of CommCell
        """
        self.cs_machine = Machine(commcell_obj.commserv_hostname, commcell_obj)
        self._cs_client = commcell_obj.clients.get(commcell_obj.commserv_name)

    def is_scheduleinfo_xml_exist(self):
        file_name = (
                self._cs_client.install_directory +
                r"\\Reports\\CommservSurvey\\ScheduleInfo.xml"
        )
        return self.cs_machine.check_file_exists(file_name)


class DocManager:
    """Provides operations on docx files"""

    def __init__(self, file):
        self.file = file
        self.fullText = []

    def read_doc(self):
        """
        Opens the document and read paragraphs
        Returns List containing text of the document
        """
        document = docx.Document(self.file)
        for para in document.paragraphs:
            self.fullText.append(para.text)
        return self.fullText

    def search_text(self, charter_data):
        """
        Validates data in the document
        """
        for each_line in charter_data:
            if each_line in self.fullText:
                continue
            else:
                raise Exception(each_line + " not found in Charter document")


class PDFManager:
    """Provides operations on a PDF File"""

    def __init__(self, pdf_file_path):
        self.pdf_file_path = pdf_file_path
        la_params = {
            'line_margin': 0.6,
        }
        self.document = load_file(self.pdf_file_path, la_params=la_params)

    def number_of_pages(self):
        return self.document.number_of_pages

    def fetch_table_data(self, page_number):
        try:
            table_data = tables.extract_table(
                self.document.elements.filter_by_page(page_number),
                as_text=True,
                fix_element_in_multiple_rows=True,
                fix_element_in_multiple_cols=True,
            )
            if not table_data:
                return "No Table found"
            return table_data

        except Exception as e:
            raise Exception(
                f"An error occurred while extracting the table : {e}"
            )
